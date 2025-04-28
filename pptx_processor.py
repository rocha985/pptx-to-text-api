from pptx import Presentation
from io import BytesIO

def extract_text_from_pptx(file_content: bytes) -> dict:
    prs = Presentation(BytesIO(file_content))
    
    slides_text = {}

    for slide_index, slide in enumerate(prs.slides, start=1):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                if slide.shapes.index(shape) == 0:
                    title = shape.text.strip()
                    slides_text[f"slide_{slide_index}"] = {"title": title, "content": []}
                else:
                    slides_text[f"slide_{slide_index}"]["content"].append(shape.text.strip())
        
        if f"slide_{slide_index}" in slides_text:
            slides_text[f"slide_{slide_index}"]["content"] = " ".join(slides_text[f"slide_{slide_index}"]["content"])

    return slides_text
